import tkinter as tk
from tkinter import filedialog
import pandas as pd
import pptx

# Create the GUI
root = tk.Tk()
root.title("RAI and TOP Meeting")
root.geometry("500x300")

# Define the function to select excel files
def select_files():
    # Open a file dialog to select the excel files
    files = filedialog.askopenfilenames(filetypes=[("Excel files", "*.xlsx")])
    return files

# Define the function to process the charts and tables
def process_files():
    # Code to process the charts and tables goes here
    pass

# Define the buttons
select_button = tk.Button(root, text="Select Files", command=select_files)
select_button.pack()

process_button = tk.Button(root, text="Process Files", command=process_files)
process_button.pack()

# Start the GUI loop
root.mainloop()

# Import required libraries
import pandas as pd
import os

# Define function to open and read Excel files
def read_excel_files(file_paths):
    charts = []
    tables = []
    for file_path in file_paths:
        if file_path.endswith('.xlsx'):
            # Open Excel file and read all sheets
            excel_file = pd.ExcelFile(file_path)
            sheets = excel_file.sheet_names
            for sheet in sheets:
                # Read each sheet and append charts and tables to corresponding lists
                df = pd.read_excel(excel_file, sheet_name=sheet)
                for i, row in df.iterrows():
                    if 'chart' in str(row).lower():
                        charts.append((file_path, sheet, i+2))
                    elif 'table' in str(row).lower():
                        tables.append((file_path, sheet, i+2))
    return charts, tables


# Import required libraries
import win32com.client as win32

# Define function to generate PowerPoint presentation with charts and tables
def generate_ppt(charts, tables, dimensions):
    # Initialize PowerPoint
    powerpoint = win32.gencache.EnsureDispatch('Powerpoint.Application')
    powerpoint.Visible = True
    presentation = powerpoint.Presentations.Add()
    
    # Add slides with charts
    for i, chart in enumerate(charts):
        file_path, sheet, row = chart
        chart_name = pd.read_excel(file_path, sheet_name=sheet, usecols='B', nrows=1, skiprows=row-1).iloc[0]
        slide = presentation.Slides.Add(i+1, 11)
        slide.Shapes.AddPicture(file_path + '!' + sheet + '$A' + str(row), 
                                LinkToFile=False, SaveWithDocument=True, 
                                Left=dimensions['chart']['left'][i], 
                                Top=dimensions['chart']['top'][i], 
                                Width=dimensions['chart']['width'], 
                                Height=dimensions['chart']['height'])
        slide.Shapes.Title.TextFrame.Text = chart_name
    
    # Add slides with tables
    for i, table in enumerate(tables):
        file_path, sheet, row = table
        table_name = pd.read_excel(file_path, sheet_name=sheet, usecols='B', nrows=1, skiprows=row-1).iloc[0]
        slide = presentation.Slides.Add(len(charts)+i+1, 11)
        table_range = pd.read_excel(file_path, sheet_name=sheet, usecols='B:G', skiprows=row-1)
        table_values = [list(table_range.columns)] + table_range.values.tolist()
        table_shape = slide.Shapes.AddTable(len(table_values), len(table_values[0]), 
                                            dimensions['table']['left'][i], 
                                            dimensions['table']['top'][i], 
                                            dimensions['table']['width'], 
                                            dimensions['table']['height'])
        for r in range(len(table_values)):
            for c in range(len(table_values[r])):
                table_shape.Table.Cell(r+1, c+1).Shape.TextFrame.Text = str(table_values[r][c])
        slide.Shapes.Title.TextFrame.Text = table_name
    
    # Save PowerPoint presentation
    presentation.SaveAs(os.getcwd() + '\\Monthly Report.pptx')
    presentation.Close()
    powerpoint.Quit()

from tkinter import messagebox

def compare_numbers(powerpoint, tables):
    # Loop through each chart and table to compare their values
    for chart_index, chart in enumerate(powerpoint.slides[0].shapes):
        if chart.has_chart:
            chart_values = chart.chart.series[0].values
            chart_title = chart.chart.chart_title.text_frame.text.lower().strip()
            for table_index, table in enumerate(tables):
                table_title = table.columns[0].value.lower().strip()
                if chart_title == table_title:
                    table_values = [row[1].value for row in table.rows]
                    for i in range(len(chart_values)):
                        if chart_values[i] != table_values[i]:
                            messagebox.showerror("Error", "Values in chart '{}' do not match values in table '{}'.".format(chart_title, table_title))
                            return False
    return True


def save_powerpoint(powerpoint):
    # Save the PowerPoint presentation to a file
    filename = filedialog.asksaveasfilename(defaultextension=".pptx")
    if filename:
        powerpoint.save(filename)
        messagebox.showinfo("Success", "PowerPoint presentation saved to '{}'.".format(filename))

